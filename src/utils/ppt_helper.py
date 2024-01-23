import math
import os
from typing import Any

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.text import MSO_VERTICAL_ANCHOR, MSO_AUTO_SIZE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
from pydantic import BaseModel


class SlideConfig(BaseModel, arbitrary_types_allowed=True):
    title: str
    x_axis_key: str
    data: pd.DataFrame
    image_path: str


class SingleChartConfig(BaseModel, arbitrary_types_allowed=True):
    index: int
    chart_type: Any
    config: SlideConfig


SERIES = ["dpo", "dso", "dio", "ccc"]


def plot_graph(config: SlideConfig):
    """
    function to create plot chart.
    """
    # Reset the plt
    df = config.data
    filter_col = config.x_axis_key
    plt.figure(figsize=(16, 7))
    if not df.empty and set([filter_col] + SERIES).issubset(df.columns):
        plt.bar(df[filter_col], df["dpo"], color="black")
        plt.bar(df[filter_col], df["dso"], color="red")
        plt.bar(df[filter_col], df["dio"], color="grey", bottom=df["dso"])
        plt.plot(df[filter_col], df["ccc"], color="black", marker="o")
    plt.ylabel("days", fontweight="bold", fontsize=15)
    plt.title(config.title)
    # plt.legend(df.columns.drop(filter_col))
    plt.legend(SERIES)
    plt.savefig(config.image_path)
    # remove the plt
    plt.clf()
    plt.cla()
    # plt.show()


def convert_df_column_to_list(df, key):
    return df[key].tolist() if not df.empty and key in df.columns else []


def create_single_chart(slide, chart_config: SingleChartConfig):
    slide_config = chart_config.config
    df = slide_config.data
    chart_width = 4
    chart_height = 3
    chart_data = CategoryChartData()
    chart_data.categories = convert_df_column_to_list(df, slide_config.x_axis_key)
    for s in SERIES:
        category_data = convert_df_column_to_list(df, s)
        if category_data:
            numeric_data = [x if isinstance(x, (int, float)) and not math.isnan(x) else 0 for x in category_data]
            chart_data.add_series(s, numeric_data)

    x, y, cx, cy = (
        Inches(1 + chart_width * (chart_config.index % 2)),
        Inches(1 + chart_height * (math.floor(chart_config.index / 2) % 2)),
        Inches(chart_width),
        Inches(chart_height),
    )
    chart = slide.shapes.add_chart(chart_config.chart_type, x, y, cx, cy, chart_data).chart
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = slide_config.title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)
    chart.chart_title.text_frame.paragraphs[0].font.bold = False

    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.size = Pt(8)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = True
    chart.legend.font.size = Pt(8)


def set_text_style(text_box, font_size=8, bold=False):
    text_box.text_frame.paragraphs[0].font.name = "Univers Next for HSBC Light"
    text_box.text_frame.paragraphs[0].font.size = Pt(font_size)
    text_box.text_frame.paragraphs[0].font.bold = bold
    text_box.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    text_box.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE


def add_title_to_slide(slide, title):
    title_shape = slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(9), Inches(1))
    title_shape.text = title
    set_text_style(title_shape, 16, True)


def create_table(df, slide, company_name):
    cols = ["dso", "dpo", "dio", "ccc"]

    # Calculate the rankings
    for col in cols:
        df[col + "_rank"] = df[col].rank()

    # Make all the numbers positive
    df.iloc[:, 1:] = df.iloc[:, 1:].abs()

    # Set Target company as the benchmark
    # benchmark = df[df["company"] == company_name].iloc[0, 1:]

    # Create the tables
    tables = []
    location = [
        [Inches(0.03), Inches(1.6), Inches(2), Inches(2)],
        [Inches(2.57), Inches(1.6), Inches(2), Inches(2)],
        [Inches(4.97), Inches(1.6), Inches(2), Inches(2)],
        [Inches(7.37), Inches(1.6), Inches(2), Inches(2)],
    ]

    for col in cols:
        table = df[["company", col, col + "_rank"]].sort_values(by=col + "_rank")
        table = table.reset_index(drop=True)
        table["rank"] = table.index + 1
        table = table[["company", col, "rank"]]
        table.index = table.index
        table = table.sort_index()

        table.loc[len(df)] = ["Peer Median", df[col].median(), ""]
        # table= pd.concat([pd.DataFrame({'company':['Dummy'], col: [0.0], 'Rank': [0]}), table])
        # table= pd.concat([pd.DataFrame({'company':['Dummy'], col: [0.0], 'Rank': [0]}), table])
        tables.append([col, table])

    txtTables = trial(tables, location, slide, company_name)

    return (txtTables, tables)


def SubElement(parent, tagname, **kwargs):
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def _set_cell_border(cell, border_color="#000000", border_width="12700"):
    """Hack function to enable the setting of border width and border color
    - top border
    - bottom border
    """
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()

    for lines in ["a:lnT", "a:lnB"]:
        # Every time before a node is inserted, the nodes with the same tag should be removed.
        tag = lines.split(":")[-1]
        for e in tcPr.getchildren():
            if tag in str(e.tag):
                tcPr.remove(e)

        ln = SubElement(tcPr, lines, w=border_width, cap="flat", cmpd="sng", algn="ctr")
        solidFill = SubElement(ln, "a:solidFill")
        SubElement(solidFill, "a:srgbClr", val=border_color)
        SubElement(ln, "a:prstDash", val="solid")
        SubElement(ln, "a:round")
        SubElement(ln, "a:headEnd", type="none", w="med", len="med")
        SubElement(ln, "a:tailEnd", type="none", w="med", len="med")

    return cell


def trial(tables, location, slide, company_name):
    txtTables = []
    tmp_i = -1
    for col, df in tables:
        tmp_i += 1
        # define the table dimensions
        rows, cols = df.shape
        rows += 1
        # define the table position and size
        left = location[tmp_i][0]
        top = location[tmp_i][1]
        width = location[tmp_i][2]
        height = round(0.3 * rows, 1)

        # add the table to the slide
        table = slide.shapes.add_table(rows, cols, left, top, width, Inches(height)).table

        # set the column widths
        table.columns[0].width = Inches(1.2)
        table.columns[1].width = Inches(0.6)
        table.columns[2].width = Inches(0.6)

        # add the table headers
        table.cell(0, 0).text = "Companies"
        table.cell(0, 1).text = col
        table.cell(0, 2).text = "Rank"

        imcd = [
            company_name,
            str(df[df["company"] == company_name][col].values[0]),
            str(df[df["company"] == company_name]["rank"].values[0]),
        ]
        df = df.drop(int(imcd[2]) - 1)

        ############################################
        # Comment this section to remove the warning
        ############################################
        # To add borders to the cells
        # for i in range(rows):
        #     for j in range(cols):
        #         cell = table.cell(i, j)
        #         cell = _set_cell_border(cell)
        ############################################
        ############################################

        for i in range(0, rows - 1):
            if i > 0:
                table.cell(i + 1, 0).text = str(df.iloc[i - 1, 0])
                table.cell(i + 1, 1).text = str(df.iloc[i - 1, 1])
                table.cell(i + 1, 2).text = str(df.iloc[i - 1, 2])

            if i == 0:
                for j in range(3):
                    table.cell(i, j).text_frame.paragraphs[0].font.bold = True
                    table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            if i < rows - 1:
                for j in range(3):
                    table.cell(i, j).fill.solid()
                    table.cell(i, j).fill.fore_color.rgb = RGBColor(255, 255, 255)  # white
                    table.cell(i, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        for j in range(3):
            table.cell(rows - 1, j).fill.solid()
            table.cell(rows - 1, j).fill.fore_color.rgb = RGBColor(175, 175, 175)  # gray
            table.cell(rows - 1, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
            table.cell(rows - 1, j).text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

        for j in range(3):
            table.cell(1, j).text = imcd[j]
            table.cell(1, j).text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 0, 0)
            table.cell(1, j).text_frame.paragraphs[0].font.bold = True

        # Setting the font size of the table to 12
        for cell in table.iter_cells():
            cell.text_frame.word_wrap = False
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)  # set the font size to 12 points

        txtTables.append(table)

        # save the presentation
    return txtTables


def create_chart_slide(slide, config: SlideConfig, title):
    add_title_to_slide(slide, title)

    create_single_chart(
        slide,
        chart_config=SingleChartConfig(index=0, config=config, chart_type=XL_CHART_TYPE.COLUMN_STACKED),
    )
    create_single_chart(slide, chart_config=SingleChartConfig(index=1, config=config, chart_type=XL_CHART_TYPE.LINE))

    # Adding 1st Image with border
    pic = slide.shapes.add_picture(config.image_path, Inches(1), Inches(4.25), width=Inches(8), height=Inches(2.5))
    pic.line.width = Pt(1)
    pic.line.color.rgb = RGBColor(225, 225, 225)

    # Adding a text box with "Source" text
    source_textbox = slide.shapes.add_textbox(Inches(1), Inches(6.75), Inches(2), Inches(0.5))
    source_textbox.text = "Source: S&P Capital IQ"
    set_text_style(source_textbox)


# End: Functions
#################################################################


def ppt_generation(configs: list[SlideConfig], company_name, llm_insight, output_path, n=6):
    """
    Function to generate the 2 ppt slide as template
    df: the benchmark table to be compare
    llm_insight: the insights generated by LLM on the company performance
    company_name: input of company name
    n: the ppt slide_layouts format, range from 0-9, 6 is a blank page

    Return: it is auto saved the ppt, so need FE to pick it up
    """
    # initial the property
    X = Presentation()

    # Create a new slide layout with a title and text box
    title_and_text_layout = X.slide_layouts[n]

    ##################################################################
    #   Start: First Slide
    first_config, second_config = configs

    if not first_config.data.empty:
        plot_graph(first_config)

        first_slide = X.slides.add_slide(title_and_text_layout)

        create_chart_slide(first_slide, first_config, "Working Capital - CCC Trend")

    df_all = second_config.data

    # Start: Second slide
    if not df_all.empty and "company" in df_all.columns:
        plot_graph(second_config)

        second_slide = X.slides.add_slide(title_and_text_layout)

        create_chart_slide(second_slide, second_config, "Working Capital - Peer Comparison")

        #   Start: Third Slide
        third_slide = X.slides.add_slide(title_and_text_layout)

        # Load the data from the Excel files
        df_all["company"] = df_all["company"].str.strip()
        df_all = df_all[df_all["company"] != "Peer Group (Median)"]
        # Fill empty cells with ""
        df_all = df_all.fillna("")

        # Putting the tables in the presentation
        create_table(df_all, third_slide, company_name)

        # Remove the "Click to add title" box <- since we changed it to a blank page so this can be remved
        # second_slide.shapes.title.text = " "

        # Add a Heading to the slide
        title_shape = third_slide.shapes.add_textbox(Inches(0.5), Inches(0), Inches(9), Inches(1))
        title_shape.text = "Relative Working Capital Efficiency Against Peer Group"
        set_text_style(title_shape, 16, True)

        # Adding a text box with "Sub-Heading" text
        source_textbox = third_slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(2), Inches(0.5))
        source_textbox.text = "Working Capital Efficiency Rankings"
        set_text_style(source_textbox, 12, True)

        table_height = round((df_all.shape[0] + 2) * 0.3, 1) + 1.6

        # Adding a text box with "Source" text
        source_textbox = third_slide.shapes.add_textbox(Inches(0.5), Inches(table_height), Inches(2), Inches(0.5))
        source_textbox.text = "Source: S&P Capital IQ"
        set_text_style(source_textbox)

        # Adding a text box with "Key Insigts" text
        source_textbox = third_slide.shapes.add_textbox(Inches(0.5), Inches(table_height + 0.5), Inches(2), Inches(0.5))
        source_textbox.text = "Key Insights"
        set_text_style(source_textbox, 12, True)

        # Adding a text box with Key Insigts information
        source_textbox = third_slide.shapes.add_textbox(Inches(0.6), Inches(table_height + 1), Inches(8), Inches(2))
        set_text_style(source_textbox)
        source_textbox.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        source_textbox.text_frame.word_wrap = True
        source_textbox.text_frame.margin_left = 0
        source_textbox.text_frame.margin_top = 0
        source_textbox.text_frame.margin_bottom = 0
        paragraph = source_textbox.text_frame.add_paragraph()
        paragraph.text = llm_insight
        paragraph.font.name = "Univers Next for HSBC Light"
        paragraph.font.size = Pt(8)

    #   End: Third Slide
    ##################################################################

    ##################################################################
    #   Save the ppt
    X.save(output_path)
    ##################################################################


def generate_ppt(df_current, df_all, insight, files_dir) -> str:
    company_name = df_current.iloc[0]["company"] if not df_current.empty and "company" in df_current.columns else ""
    if company_name:
        latest_days = df_all.iloc[0]["days"] if not df_all.empty and "days" in df_all.columns else ""
        df_all = df_all[df_all["days"] == latest_days] if latest_days else pd.DataFrame([])
        os.makedirs(files_dir, exist_ok=True)
        first_slide_config = SlideConfig(
            image_path=f"{files_dir}/image1.png",
            title=f"{company_name} – Working Capital Metrics Trend",
            x_axis_key="days",
            data=df_current.drop(columns=["company"], errors="ignore"),
        )
        second_slide_config = SlideConfig(
            image_path=f"{files_dir}/image2.png",
            title=f"Peer Comparison ({latest_days})",
            x_axis_key="company",
            data=df_all.drop(columns=["days"], errors="ignore"),
        )
        ppt_path = f"{files_dir}/ppt.pptx"
        ppt_generation(
            configs=[first_slide_config, second_slide_config],
            company_name=company_name,
            llm_insight=insight,
            output_path=ppt_path,
        )
        return ppt_path
    return ""


if __name__ == "__main__":
    test_formatted_data = {
        "all_years_data": [
            {"company": "CDT N.V.", "days": "FY18", "dpo": -45, "dso": 53, "dio": 61, "ccc": 69},
            {"company": "CDT N.V.", "days": "FY19", "dpo": -47, "dso": 53, "dio": 64, "ccc": 70},
            {"company": "CDT N.V.", "days": "FY20", "dpo": -49, "dso": 56, "dio": 64, "ccc": 71},
            {"company": "CDT N.V.", "days": "FY21", "dpo": -46, "dso": 55, "dio": 63, "ccc": 72},
        ],
        "latest_year_data": [
            {"company": "ABC Holding", "days": "FY21", "dpo": -78, "dso": 67, "dio": 41, "ccc": 30},
            {"company": "Uni Inc.", "days": "FY21", "dpo": -44, "dso": 53, "dio": 41, "ccc": 50},
            {"company": "DEB Company", "days": "FY21", "dpo": -39, "dso": 49, "dio": 43, "ccc": 53},
            {"company": "CCT N.V.", "days": "FY21", "dpo": -46, "dso": 38, "dio": 61, "ccc": 53},
            {"company": "BZL plc", "days": "FY21", "dpo": -54, "dso": 41, "dio": 68, "ccc": 55},
            {"company": "CDT N.V.", "days": "FY21", "dpo": -46, "dso": 55, "dio": 63, "ccc": 72},
        ],
    }
    current_data_test = test_formatted_data["all_years_data"] if "all_years_data" in test_formatted_data else []
    all_data_test = test_formatted_data["latest_year_data"] if "latest_year_data" in test_formatted_data else []
    df_current_test = pd.DataFrame(current_data_test).drop(
        columns=["meta__score", "meta__reference", "id"], errors="ignore"
    )
    df_all_test = pd.DataFrame(all_data_test).drop(columns=["meta__score", "meta__reference", "id"], errors="ignore")
    test_insight = """
    Q1: Apple Inc.'s CCC increased from 53 to 72 days (FY13-FY21) due to higher DSO and DIO. Reasons could be slower collections and inventory management.
    Q2: Apple Inc.'s CCC (72) is higher than the peer group median (57), indicating less efficiency in working capital management.
    """
    test_dir = os.path.join(os.path.dirname(__file__), "../../", "tmp/wcs")
    os.makedirs(test_dir, exist_ok=True)
    generate_ppt(df_current_test, df_all_test, test_insight, test_dir)
